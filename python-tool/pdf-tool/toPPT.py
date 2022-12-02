import fitz
import pptx
from pptx.util import Inches

'''
    本代码功能实现
        PDF转成图片
        图片转成PPT
'''

# 图片转PPT
def img2pptx(img_path, ppt_name):
    ppt = pptx.Presentation()  # 生成ppt对象
    for i in img_path:
        layout = ppt.slide_layouts[6]  # 定义一个 PPT 页面 插入图片,这里选择样式6
        slide = ppt.slides.add_slide(layout)

        image = slide.shapes.add_picture(
            image_file=i,
            left=Inches(0),
            top=Inches(0),
            width=Inches(10),
            height=Inches(8)
        )
        # 目前设置会撑满整张PPT，大小可自行调整
    ppt.save(ppt_name)
    print('完成了！')

    # PDF转图片
def pdf2img(pdf_path, img_dir, ppt_name):
    doc = fitz.open(pdf_path)  # 打开pdf
    img_path = []
    for page in doc:  # 遍历pdf的每一页
        zoom_x = 2.0  # 设置每页的水平缩放因子
        zoom_y = 2.0  # 设置每页的垂直缩放因子
        mat = fitz.Matrix(zoom_x, zoom_y)
        pix = page.get_pixmap(matrix=mat)
        pix.save(r"{}page-{}.png".format(img_dir, page.number))  # 保存
        img_path.append(r"{}page-{}.png".format(img_dir, page.number))

    img2pptx(img_path, ppt_name)


if __name__ == '__main__':

    print('开始了，请等待……')
    # pdf路径
    pdf_path = "pdf/jianzhu.pdf"
    # 图片保存位置
    img_dir = "img/"

    # ppt名称
    ppt_name = "ppt/result.pptx"

    # pdf转图片
    pdf2img(pdf_path, img_dir, ppt_name)
